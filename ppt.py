from PIL import Image
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('hebe.pptx')
def handle_sld(i):
	sld = prs.slides[i]
	shapes = sld.shapes
	index = i + 1

	name_shapes = [s for s in shapes if s.shape_type in [MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.AUTO_SHAPE] and s.has_text_frame and s.text != '']
	if len(name_shapes) != 1:
		print(index, 'error1')
		return
	name = name_shapes[0].text
	pic_n = 1
	pics = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(s, 'image')]
	if len(pics) == 0:
		print(index, 'error2')
		return
	for pic in pics:
		imageStream = io.BytesIO(pic.image.blob)
		imageFile = Image.open(imageStream)
		img_name = f'{index}-{name}-{pic_n}.png'
		imageFile.save(f'hebe-avatar/{img_name}')
		pic_n += 1
	print(index, 'success')

for i in range(69):
	handle_sld(i)